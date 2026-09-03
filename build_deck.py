"""Build the 10-slide Goa Finals deck from the shipped run's own figures.

Every image on a slide is a file in `figures/`, written by `figures.py` from the delivered
CSVs. Nothing here re-derives a number and nothing here is drawn by hand, so a slide cannot
disagree with the run the way a hand-built deck eventually does.

Each slide carries a title, a one-line kicker, the figure, and speaker notes sized for a
ten-minute talk. The notes are the script; the slide is what the room reads while it is
being said.

    python build_deck.py        # writes Sokhda_Goa_Finals.pptx
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.abspath(__file__))
FIGURES = os.path.join(ROOT, "figures")
OUT = os.path.join(ROOT, "Sokhda_Goa_Finals.pptx")

W, H = Inches(13.333), Inches(7.5)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x5F, 0x6B)

SLIDES = [
    dict(
        figure="cover",
        title="Sokhda kharif 2025 — final yield forecast",
        kicker="966 plots · six Capella X-band passes · no ground truth · "
               "893.9 t over 447.5 ha, 2.00 t/ha",
        notes="""\
One village, Sokhda in Vadodara district. 966 plots, median a quarter of a hectare. Six
Capella X-band HH passes, June to November. No labels, no leaderboard.

Two things the brief says that the files do not. It says an expanded set of villages; the
village shapefile holds one polygon. It says the crop classification is carried forward; the
farm shapefile has five fields and none is a crop. So we re-derived the labels.

There is nothing to fit, so validation is the deliverable. Seventeen predictions were written
down before the data that could test them was opened, and nine were contradicted.""",
    ),
    dict(
        figure="sar_composite",
        title="What the radar actually sees",
        kicker="Three dates in three colour channels. No model, no classifier, no fitted "
               "number.",
        notes="""\
Before any method, this. June in red, August in green, November in blue -- the same
calibrated gamma-nought, nothing else done to it. The field pattern falls out of the colour.

Green is a canopy that peaked in August and was gone by November. Blue and magenta are still
bright on the twelfth -- cotton, and the twelve long-duration parcels. Grey did not change:
roads, bunds, the village. If X-band could not see this season plot by plot, it would be grey
throughout.""",
    ),
    dict(
        figure="trajectories",
        title="The stack, and the pass that fights you",
        kicker="T5 is right-looking, pre-dawn, and follows 63 mm of rain. Its level is "
               "never used.",
        notes="""\
Five passes look left from about 135 degrees. T5, the 29th of October, looks right from 318 --
shadow and layover fall on the opposite side of every bund. It is also a 01:37 pass at maximum
canopy dew, following 63 millimetres of rain in three days.

It broke the co-registration matcher, reporting a 108 metre shift. We diagnosed instead of
clamping: its correlation peak is a third of the stack's and nearly flat, the reversed look
expressed as a statistic.

Its residual is not a constant offset -- it changes sign with brightness -- so T5's level is
replaced by the T4-T6 interpolation. Only the residual survives, as a weak covariate.""",
    ),
    dict(
        figure="canopy_sign",
        title="We pre-registered the canopy sign. We were wrong.",
        kicker="Predicted attenuation on four of five crops. Measured +1 on all five, "
               "rho = +0.569, n = 813.",
        notes="""\
Four of the five crops are darkest at what looks like peak canopy, which reads as the canopy
attenuating the surface return. We wrote that into a module constant above the code that opens
the optical file, and it has not been edited since.

Sentinel-2 lands on a Capella date twice. We difference both instruments, so plot size, soil
texture and row orientation cancel. The answer is positive on all five crops. Greener is
brighter here. Our sign-agnostic design scored minus 0.085; the signed form scores plus 0.564.
It was not conservative, it was empty, and we rebuilt the module rather than patch it.

And where that arbiter came from: luck. Cloud-free optical coincided with a Capella pass twice
in six, and T5's only candidate was 79 percent cloud, which is why the T5 control does not
exist. The one thing the radar could not settle alone needed a second instrument on the same
day, and we got one by coincidence rather than design.

The same test refused the per-plot harvest date we promised, and we deleted it: three canopy
observations with a sixty-day gap cannot locate a transition.""",
    ),
    dict(
        figure="crop_type_map",
        title="Crop labels, re-derived from six dates",
        kicker="Tier-1 coverage fell to 26.5 % — the target was missed, and stability "
               "roughly doubled.",
        notes="""\
Two November acquisitions add the discriminator four dates could not: bajra off the field by
late September, maize in October, groundnut lifted October-November, cotton still standing.
Cotton does not cluster, so its rule is per plot and in absolute decibels.

We aimed to raise tier-1 coverage above Round 2's 31.6 percent. We did not: 26.5, a missed
target. What improved is stability -- tier-1 is 100 percent stable across every clustering
setting, where Round 2's could halve.

Two defects, one found by Kaggle and one by an audit of ourselves. The ranking axis for the
allocated remainder was the November canopy clipped at zero, so 403 of 793 plots sat at exactly
zero and sort order decided the bajra-maize cut. Two machines, same code, 39 plots different.
The signed departure separates that block across 392 values.

We then reported the re-ranking confirmed by optical data. It is not. Our own audit found the
test still residualising against the axis from before the fix. Corrected, tier-2 labels carry
no optical information beyond their own axis while the tier-1 control does. Tier 2 is an
allocation, which is what this slide has called it all along.""",
    ),
    dict(
        figure="model_chain",
        title="One measured term on a sourced reference",
        kicker="Y_final = Y_ref(crop, 2025) × a(season-complete canopy integral)",
        notes="""\
One modulation term, not three. Round 2 measured its own problem: within a cohort its health
index and yield estimate correlated at exactly 1.000 -- one ranking under two names.

The reference yield inverted our planning assumption. Sokhda's monsoon was 119 percent of its
thirty-year mean, so we had planned to adjust last year's yield upward. The state estimate
says the opposite: Gujarat kharif rice and bajra hit five-year lows, minus 29 and minus 26
percent, in an excess-rain season -- the Narmada overflowed inside the paddy grain-fill
window. Had we applied our elasticity, rice would have been forecast above a reference the
state measured 29 percent below. That correction came from checking a source, not modelling.""",
    ),
    dict(
        figure="yield_forecast_map",
        title="The forecast, plot by plot",
        kicker="893.9 t · 447.5 ha · 2.00 t/ha · ±151 t of that is external assumption, "
               "±9.5 t is the radar",
        notes="""\
Groundnut 341 plots and 332 tonnes, maize 313 and 274, rice 111 and 128, bajra 139 and 87,
cotton 62 and 74. The shipped table is 966 rows and 21 columns and carries the whole chain,
not the answer alone.

We priced the total by re-running the chain under each source. State reference at a stated ten
percent, plus or minus 89 tonnes. The district crop mix, which allocates three of our five
cohorts, 62. Round 2's labels 7, speckle 3, tie ordering zero.

The mix is the row worth explaining, because we can score it against itself. Rice and cotton
are assigned by threshold rules, not by the mix. The district says rice is 26 percent of area
and cotton 32; we measure 17 and 10. It overstates both, and that disagreement is what we use
to perturb the three crops we cannot check.

External assumptions, 151 tonnes. Everything the radar contributes, 9.5. Somebody else's
numbers decide where the line is and how big three of the five cohorts are.""",
    ),
    dict(
        figure="extrapolation",
        title="What makes it a forecast and not a restatement",
        kicker="Four crops closed by observation. Cotton alone is 56 % projection, and it "
               "says so.",
        notes="""\
Round 2 discounted the unobserved rest of the season with a hand-set constant. Round 3
replaces the discount with measurement, because the stack now contains the harvest for four
of the five crops.

Cotton is the only crop whose season runs past the 12th of November, and the only one with a
projected share: 56 percent of its canopy-days. That ships per plot.

Nothing in our data observes that window, so we went and found something that does. Sixteen
free Sentinel-1 passes, feeding no feature and no label. Cotton is the only cohort still above
its own June bare soil after the 12th of November, and it rises nearly a decibel through to
the 21st of December. The flat hold is not optimistic. That was written down before we looked.

The projection is flat -- last observed canopy carried forward. That is not the rule we
started with, and the next slide is why.""",
    ),
    dict(
        figure="backtest",
        title="The back-test deleted our own rule",
        kicker="Shipped rule vs persistence: −0.119, 95 % interval [−0.280, +0.022]. It "
               "does not beat persistence.",
        notes="""\
Fit on T1 to T4, predict the withheld 12th of November pass, on Round 2's four-date labels so
no November information leaks.

Our decaying projection first scored plus 0.284 against persistence. We quote that nowhere. It
predicts a higher canopy, and the scene sits 1.65 decibels above June, so it could be winning
by offsetting a drift neither predictor models. Hand every predictor that drift and it scores
minus 0.409. We deleted it.

The flat hold we shipped scores minus 0.119, interval containing zero -- and minus 0.216 on
the 732 plots where the rule actually fires, which is the fairer number and the worse one.

The back-test's value was not certifying the model. It was deleting a rule that looked
principled, gave a favourable headline, and did not survive a control built to break it.""",
    ),
    dict(
        figure="reserved_optical",
        title="Two scenes nothing upstream was allowed to read",
        kicker="Cotton's December NDVI 0.690 vs 0.474–0.532, one-sided p = 1.26e-11.",
        notes="""\
The 12th of December and the 16th of January were reserved from the first fetch, and an
assertion greps the source tree and fails the run if any module but the validator names them.
That is a lint, not a proof, and our own audit says so. Both dates sit inside the rabi window,
so what they test is which plots still carry a kharif crop after everything else has finished
-- and cotton alone is picked into January.

Cotton's December NDVI is 0.690 against 0.474 to 0.532. A SAR-only label picked the right
plots on a scene it never saw. The negative control matters as much: cleared plots are under
rabi, not bare, so this is not soil quality.

One more falsification. Twelve parcels we called an orchard were confirmed in December and
January -- and refused in June, 0.247 against 0.397. Bare in June, so a long-duration monsoon
crop. We renamed the constant and kept the finding.""",
    ),
    dict(
        figure="zone_map",
        title="Aggregation, and what we do not claim",
        kicker="46 zones, 1.50–2.80 t/ha around a village 2.00. Relabelling moves the "
               "total by 1.5 %.",
        notes="""\
One village makes the required village table a single row. The 500 metre grid is what makes
the aggregation an aggregation: 46 cells of at least five farms, 946 of the 966 plots,
spreading 1.50 to 2.80 tonnes per hectare around a village figure of 2.00. The village row is
the sum of the shipped plot file, rounded once before aggregating -- our own cross-check
caught a 0.0015 tonne discrepancy.

The roll-up is gated on the village geometry, not the village name: 962 of 966 plots agree
with the attribute, none disagree, and all 447 hectares sit inside a boundary enclosing 1174.
So we report 38 percent of Sokhda -- its digitised farmland, not the village.

What we do not claim. Tier-1 covers a quarter of the area, below Round 2, a missed target.
Median peak canopy is 0.77 decibels: real, corroborated, small. X-band is supposed to saturate
early, and that is the first question we expect from this room -- across six NDVI bins the
departure rises monotonically, so over the range these fields occupy it does not saturate.
Plot-level irrigation could produce the same green-and-bright correlation as canopy scattering.
And our projection is no better than persistence.

Thank you.""",
    ),
]


def _text(slide, left, top, width, height, text, size, bold=False, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Helvetica Neue"
    return box


def _fit(path: str, box_l: int, box_t: int, box_w: int, box_h: int) -> tuple:
    """Largest placement of the image inside the box that preserves its aspect ratio."""
    from PIL import Image

    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    return box_l + (box_w - w) // 2, box_t + (box_h - h) // 2, w, h


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]

    for spec in SLIDES:
        slide = prs.slides.add_slide(blank)
        _text(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7),
              spec["title"], 30, bold=True)
        _text(slide, Inches(0.6), Inches(1.02), Inches(12.1), Inches(0.5),
              spec["kicker"], 15, color=MUTED)

        path = os.path.join(FIGURES, spec["figure"] + ".png")
        if not os.path.isfile(path):
            raise FileNotFoundError(f"{path} -- run the pipeline before building the deck")
        left, top, width, height = _fit(
            path, Emu(Inches(0.6)), Emu(Inches(1.65)),
            Emu(Inches(12.1)), Emu(Inches(5.5)))
        slide.shapes.add_picture(path, left, top, width, height)

        slide.notes_slide.notes_text_frame.text = spec["notes"]

    return prs


if __name__ == "__main__":
    prs = build()
    prs.save(OUT)
    words = sum(len(s["notes"].split()) for s in SLIDES)
    print(f"wrote {OUT}: {len(SLIDES)} slides, {words} words of speaker notes "
          f"({words / 150:.1f} min at 150 wpm, {words / 140:.1f} min at 140) "
          f"against a 10-minute slot")
